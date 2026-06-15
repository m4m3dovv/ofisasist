from __future__ import annotations

import re
import shutil
from dataclasses import dataclass
from pathlib import Path

import pandas as pd
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from local_ai import ask_local_ai


TABLE_EXTENSIONS = {".xlsx", ".xls", ".csv"}
DOCUMENT_EXTENSIONS = {".docx", ".txt", ".md"}
PRESENTATION_EXTENSIONS = {".pptx"}
PDF_EXTENSIONS = {".pdf"}
SUPPORTED_EXTENSIONS = TABLE_EXTENSIONS | DOCUMENT_EXTENSIONS | PRESENTATION_EXTENSIONS | PDF_EXTENSIONS


@dataclass
class TaskResult:
    message: str
    output_path: Path | None = None


def supported_file_types() -> str:
    return ", ".join(sorted(SUPPORTED_EXTENSIONS))


def is_supported_file(path: Path) -> bool:
    return path.suffix.lower() in SUPPORTED_EXTENSIONS


def run_office_task(file_path: Path, instruction: str, output_dir: Path) -> TaskResult:
    instruction_normalized = _normalize(instruction)
    output_dir.mkdir(parents=True, exist_ok=True)

    if not is_supported_file(file_path):
        return TaskResult(f"Bu fayl tipi hələ dəstəklənmir. Hazırda bunları qəbul edirəm: `{supported_file_types()}`.")

    suffix = file_path.suffix.lower()
    if suffix in TABLE_EXTENSIONS:
        return _run_table_task(file_path, instruction, instruction_normalized, output_dir)
    if suffix in DOCUMENT_EXTENSIONS:
        return _run_text_document_task(file_path, instruction, instruction_normalized, output_dir)
    if suffix in PRESENTATION_EXTENSIONS:
        return _run_presentation_task(file_path, instruction, instruction_normalized, output_dir)
    if suffix in PDF_EXTENSIONS:
        return _run_pdf_task(file_path, instruction, instruction_normalized, output_dir)

    return TaskResult(f"Bu fayl tipi hələ dəstəklənmir. Hazırda bunları qəbul edirəm: `{supported_file_types()}`.")


def _run_table_task(
    file_path: Path,
    instruction: str,
    instruction_normalized: str,
    output_dir: Path,
) -> TaskResult:
    if _wants_columns(instruction_normalized):
        df = _read_table(file_path)
        columns = "\n".join(f"- {column}" for column in df.columns)
        return TaskResult(f"Fayldakı sütunlar:\n{columns}")

    if _wants_summary(instruction_normalized):
        return TaskResult(_summarize_table(file_path))

    if _wants_excel(instruction_normalized):
        df = _read_table(file_path)
        output_path = output_dir / f"{file_path.stem}_converted.xlsx"
        df.to_excel(output_path, index=False)
        return TaskResult("Faylı Excel formatına çevirdim.", output_path)

    if _wants_csv(instruction_normalized):
        df = _read_table(file_path)
        output_path = output_dir / f"{file_path.stem}_converted.csv"
        df.to_csv(output_path, index=False)
        return TaskResult("Faylı CSV formatına çevirdim.", output_path)

    search_result = _try_search_text(_table_to_text(file_path), instruction, instruction_normalized)
    if search_result:
        return search_result

    filter_result = _try_filter(file_path, instruction, instruction_normalized, output_dir)
    if filter_result:
        return filter_result

    sort_result = _try_sort(file_path, instruction, instruction_normalized, output_dir)
    if sort_result:
        return sort_result

    return _try_local_ai_or_help(file_path, instruction, _table_to_text(file_path))


def _run_text_document_task(
    file_path: Path,
    instruction: str,
    instruction_normalized: str,
    output_dir: Path,
) -> TaskResult:
    text = _extract_document_text(file_path)

    if _wants_txt(instruction_normalized) or _wants_extract_text(instruction_normalized):
        output_path = output_dir / f"{file_path.stem}_text.txt"
        output_path.write_text(text, encoding="utf-8")
        return TaskResult("Sənəddəki mətni `.txt` faylına çıxardım.", output_path)

    search_result = _try_search_text(text, instruction, instruction_normalized)
    if search_result:
        return search_result

    if _wants_summary(instruction_normalized):
        return TaskResult(_summarize_text_file(file_path, text))

    return _try_local_ai_or_help(file_path, instruction, text)


def _run_presentation_task(
    file_path: Path,
    instruction: str,
    instruction_normalized: str,
    output_dir: Path,
) -> TaskResult:
    slides = _extract_presentation_slides(file_path)
    text = "\n\n".join(f"Slayd {index}: {slide}" for index, slide in slides)

    if _wants_txt(instruction_normalized) or _wants_extract_text(instruction_normalized):
        output_path = output_dir / f"{file_path.stem}_slides_text.txt"
        output_path.write_text(text, encoding="utf-8")
        return TaskResult("Prezentasiyadakı mətni `.txt` faylına çıxardım.", output_path)

    search_result = _try_search_text(text, instruction, instruction_normalized)
    if search_result:
        return search_result

    if _wants_summary(instruction_normalized):
        return TaskResult(_summarize_presentation(file_path, slides))

    return _try_local_ai_or_help(file_path, instruction, text)


def _run_pdf_task(
    file_path: Path,
    instruction: str,
    instruction_normalized: str,
    output_dir: Path,
) -> TaskResult:
    pages = _extract_pdf_pages(file_path)
    if _pages_are_empty(pages):
        pages = _extract_pdf_pages_with_ocr(file_path)
    text = "\n\n".join(f"Səhifə {index}: {page}" for index, page in pages)

    if _wants_txt(instruction_normalized) or _wants_extract_text(instruction_normalized):
        output_path = output_dir / f"{file_path.stem}_pdf_text.txt"
        output_path.write_text(text, encoding="utf-8")
        return TaskResult("PDF-dəki mətni `.txt` faylına çıxardım.", output_path)

    search_result = _try_search_text(text, instruction, instruction_normalized)
    if search_result:
        return search_result

    if _wants_summary(instruction_normalized):
        return TaskResult(_summarize_pages(file_path, pages, "Səhifə sayı"))

    return _try_local_ai_or_help(file_path, instruction, text)


def _read_table(file_path: Path) -> pd.DataFrame:
    suffix = file_path.suffix.lower()
    if suffix == ".csv":
        return pd.read_csv(file_path)
    return pd.read_excel(file_path)


def _summarize_table(file_path: Path) -> str:
    df = _read_table(file_path)
    lines = [
        f"Fayl: `{file_path.name}`",
        f"Tip: cədvəl",
        f"Sətir sayı: {len(df)}",
        f"Sütun sayı: {len(df.columns)}",
        "",
        "Sütunlar:",
    ]
    lines.extend(f"- {column}" for column in df.columns)

    numeric_columns = df.select_dtypes(include="number").columns
    if len(numeric_columns) > 0:
        lines.append("")
        lines.append("Rəqəmsal sütunların qısa statistikası:")
        stats = df[numeric_columns].describe().round(2).transpose()
        for column, row in stats.iterrows():
            lines.append(f"- {column}: min={row['min']}, max={row['max']}, orta={row['mean']}")

    missing = df.isna().sum()
    missing = missing[missing > 0]
    if not missing.empty:
        lines.append("")
        lines.append("Boş xanalar:")
        lines.extend(f"- {column}: {count}" for column, count in missing.items())

    return "\n".join(lines)


def _extract_document_text(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".docx":
        document = Document(file_path)
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]

        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))

        return "\n".join(paragraphs)

    return file_path.read_text(encoding="utf-8", errors="ignore")


def _extract_presentation_slides(file_path: Path) -> list[tuple[int, str]]:
    presentation = Presentation(file_path)
    slides: list[tuple[int, str]] = []

    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        slides.append((index, "\n".join(parts)))

    return slides


def _extract_pdf_pages(file_path: Path) -> list[tuple[int, str]]:
    reader = PdfReader(str(file_path))
    pages: list[tuple[int, str]] = []

    for index, page in enumerate(reader.pages, start=1):
        pages.append((index, page.extract_text() or ""))

    return pages


def _extract_pdf_pages_with_ocr(file_path: Path) -> list[tuple[int, str]]:
    if shutil.which("tesseract") is None:
        return []

    try:
        import fitz
        import pytesseract
        from PIL import Image
    except ImportError:
        return []

    pages: list[tuple[int, str]] = []
    document = fitz.open(file_path)
    for index, page in enumerate(document, start=1):
        pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
        text = pytesseract.image_to_string(image, lang="aze+eng")
        pages.append((index, text.strip()))

    return pages


def _summarize_text_file(file_path: Path, text: str) -> str:
    paragraphs = [line for line in text.splitlines() if line.strip()]
    preview = _preview_text(text)
    return "\n".join(
        [
            f"Fayl: `{file_path.name}`",
            "Tip: sənəd",
            f"Sətir/abzas sayı: {len(paragraphs)}",
            f"Söz sayı: {_word_count(text)}",
            "",
            "İlk hissə:",
            preview or "Mətn tapılmadı.",
        ]
    )


def _summarize_presentation(file_path: Path, slides: list[tuple[int, str]]) -> str:
    non_empty_slides = [(index, text) for index, text in slides if text.strip()]
    lines = [
        f"Fayl: `{file_path.name}`",
        "Tip: prezentasiya",
        f"Slayd sayı: {len(slides)}",
        f"Mətn olan slayd sayı: {len(non_empty_slides)}",
        "",
        "Slaydlar:",
    ]

    for index, text in non_empty_slides[:10]:
        lines.append(f"- Slayd {index}: {_single_line_preview(text)}")

    if len(non_empty_slides) > 10:
        lines.append(f"- ... daha {len(non_empty_slides) - 10} slayd")

    return "\n".join(lines)


def _summarize_pages(file_path: Path, pages: list[tuple[int, str]], count_label: str) -> str:
    text = "\n".join(page_text for _, page_text in pages)
    lines = [
        f"Fayl: `{file_path.name}`",
        "Tip: PDF",
        f"{count_label}: {len(pages)}",
        f"Söz sayı: {_word_count(text)}",
        "",
        "İlk hissə:",
        _preview_text(text) or "Mətn tapılmadı. PDF skan şəklində ola bilər.",
    ]
    return "\n".join(lines)


def _table_to_text(file_path: Path) -> str:
    df = _read_table(file_path)
    return df.astype(str).to_string(index=False)


def _try_sort(
    file_path: Path,
    instruction: str,
    instruction_normalized: str,
    output_dir: Path,
) -> TaskResult | None:
    if "sirala" not in instruction_normalized and "sort" not in instruction_normalized:
        return None

    df = _read_table(file_path)
    column = _find_column(df, instruction)
    if column is None:
        return TaskResult("Sıralamaq üçün sütunu tapa bilmədim. Məsələn: `məbləğ sütununa görə sırala`.")

    descending = any(word in instruction_normalized for word in ["azalan", "boyukden", "desc"])
    sorted_df = df.sort_values(by=column, ascending=not descending)
    output_path = output_dir / f"{file_path.stem}_sorted.xlsx"
    sorted_df.to_excel(output_path, index=False)
    return TaskResult(f"`{column}` sütununa görə sıraladım.", output_path)


def _try_filter(
    file_path: Path,
    instruction: str,
    instruction_normalized: str,
    output_dir: Path,
) -> TaskResult | None:
    if "filtr" not in instruction_normalized and "filter" not in instruction_normalized:
        return None

    df = _read_table(file_path)
    match = re.search(r"(.+?)\s*(=|==|:)\s*(.+?)(?:\s+olanlari|\s+filtr|$)", instruction, flags=re.I)
    if not match:
        return TaskResult("Filtr şərtini tapa bilmədim. Məsələn: `status = ödənilib filtr et`.")

    raw_column, _, raw_value = match.groups()
    column = _find_column(df, raw_column)
    if column is None:
        return TaskResult("Filtr üçün sütunu tapa bilmədim. Sütun adını fayldakı kimi yazmağa çalışın.")

    value = raw_value.strip().strip("'\"")
    filtered = df[df[column].astype(str).str.casefold() == value.casefold()]
    output_path = output_dir / f"{file_path.stem}_filtered.xlsx"
    filtered.to_excel(output_path, index=False)
    return TaskResult(f"`{column}` = `{value}` olan {len(filtered)} sətir tapdım.", output_path)


def _try_search_text(text: str, instruction: str, instruction_normalized: str) -> TaskResult | None:
    if "axtar" not in instruction_normalized and "tap" not in instruction_normalized and "search" not in instruction_normalized:
        return None

    query = _extract_search_query(instruction)
    if not query:
        return TaskResult("Axtarılacaq sözü tapa bilmədim. Məsələn: `müştəri sözünü axtar`.")

    matches = []
    query_normalized = _normalize(query)
    for line_number, line in enumerate(text.splitlines(), start=1):
        if query_normalized in _normalize(line):
            matches.append((line_number, line.strip()))

    if not matches:
        return TaskResult(f"`{query}` üçün nəticə tapılmadı.")

    lines = [f"`{query}` üçün {len(matches)} nəticə tapdım:"]
    for line_number, line in matches[:20]:
        lines.append(f"- Sətir {line_number}: {_single_line_preview(line, 160)}")
    if len(matches) > 20:
        lines.append(f"- ... daha {len(matches) - 20} nəticə")

    return TaskResult("\n".join(lines))


def _try_local_ai_or_help(file_path: Path, instruction: str, file_text: str) -> TaskResult:
    ai_result = ask_local_ai(instruction, file_path.name, file_text)
    if ai_result.ok:
        return TaskResult(ai_result.message)

    help_message = _unknown_task_message(file_path).message
    return TaskResult(f"{help_message}\n\nLokal AI qeydi: {ai_result.message}")


def _extract_search_query(instruction: str) -> str:
    cleaned = re.sub(r"\b(axtar|tap|search|sozunu|sözünü|metnini|mətnini)\b", " ", instruction, flags=re.I)
    cleaned = cleaned.replace("`", "").replace('"', "").replace("'", "")
    return " ".join(cleaned.split()).strip()


def _find_column(df: pd.DataFrame, text: str) -> str | None:
    text_normalized = _normalize(text)
    normalized_columns = {_normalize(str(column)): str(column) for column in df.columns}

    for normalized, original in normalized_columns.items():
        if normalized and normalized in text_normalized:
            return original

    for normalized, original in normalized_columns.items():
        if normalized and text_normalized in normalized:
            return original

    return None


def _unknown_task_message(file_path: Path) -> TaskResult:
    suffix = file_path.suffix.lower()
    if suffix in TABLE_EXTENSIONS:
        examples = "`xülasə ver`, `sütunları göstər`, `status = ödənilib filtr et`, `məbləğ sütununa görə sırala`, `csv et`"
    elif suffix in PRESENTATION_EXTENSIONS:
        examples = "`xülasə ver`, `mətni çıxart`, `müştəri sözünü axtar`"
    elif suffix in PDF_EXTENSIONS:
        examples = "`xülasə ver`, `mətni çıxart`, `müqavilə sözünü axtar`"
    else:
        examples = "`xülasə ver`, `mətni çıxart`, `müştəri sözünü axtar`"

    return TaskResult(
        "Tapşırığı hazır qaydalarla tam başa düşmədim. "
        f"Bu fayl üçün belə yaza bilərsiniz: {examples}. "
        "Daha sərbəst tapşırıqlar üçün Ollama lokal AI aktiv olmalıdır."
    )


def _pages_are_empty(pages: list[tuple[int, str]]) -> bool:
    return _word_count("\n".join(text for _, text in pages)) == 0


def _wants_summary(instruction: str) -> bool:
    return any(word in instruction for word in ["xulase", "summary", "analiz", "hesabat", "melumat"])


def _wants_columns(instruction: str) -> bool:
    return any(phrase in instruction for phrase in ["sutun", "column", "basliq"])


def _wants_csv(instruction: str) -> bool:
    return "csv" in instruction


def _wants_excel(instruction: str) -> bool:
    return any(word in instruction for word in ["excel", "xlsx"])


def _wants_txt(instruction: str) -> bool:
    return any(word in instruction for word in ["txt", "text", "metn", "metin"])


def _wants_extract_text(instruction: str) -> bool:
    return any(phrase in instruction for phrase in ["metni cixart", "metn cixart", "extract"])


def _word_count(text: str) -> int:
    return len(re.findall(r"\w+", text, flags=re.UNICODE))


def _preview_text(text: str, limit: int = 900) -> str:
    compact = " ".join(text.split())
    if len(compact) <= limit:
        return compact
    return compact[: limit - 3].rstrip() + "..."


def _single_line_preview(text: str, limit: int = 120) -> str:
    compact = " ".join(text.split())
    if len(compact) <= limit:
        return compact
    return compact[: limit - 3].rstrip() + "..."


def _normalize(text: str) -> str:
    replacements = {
        "ə": "e",
        "ö": "o",
        "ü": "u",
        "ı": "i",
        "ğ": "g",
        "ş": "s",
        "ç": "c",
        "Ə": "e",
        "Ö": "o",
        "Ü": "u",
        "I": "i",
        "İ": "i",
        "Ğ": "g",
        "Ş": "s",
        "Ç": "c",
    }
    lowered = text.casefold()
    for original, replacement in replacements.items():
        lowered = lowered.replace(original, replacement)
    return lowered
